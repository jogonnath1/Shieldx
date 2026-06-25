import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
# Set to Widescreen 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def apply_template(slide, slide_number):
    # Background (Solid light blue-gray)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(226, 239, 244)
    
    # Left dark teal strip
    left_strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), Inches(7.5)
    )
    left_strip.fill.solid()
    left_strip.fill.fore_color.rgb = RGBColor(47, 78, 90)
    left_strip.line.fill.background()
    
    # Optional decorative curved lines (using arc shapes as a hack)
    arc = slide.shapes.add_shape(
        MSO_SHAPE.ARC, Inches(-2), Inches(1), Inches(4), Inches(6)
    )
    arc.line.color.rgb = RGBColor(47, 78, 90)
    arc.line.width = Pt(1)
    
    # Slide number shape (Chevron)
    try:
        slide_num_shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, 0, Inches(5.0), Inches(1.8), Inches(0.8)
        )
    except:
        slide_num_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(5.0), Inches(1.8), Inches(0.8)
        )
    slide_num_shape.fill.solid()
    slide_num_shape.fill.fore_color.rgb = RGBColor(58, 58, 58)
    slide_num_shape.line.fill.background()
    
    # Add text to it
    tf = slide_num_shape.text_frame
    tf.text = str(slide_number)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Times New Roman"
    # To center in the chevron, add spaces or margin (hack)
    tf.margin_left = Inches(0.8)

def add_styled_slide(title, bullet_points, slide_number):
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    apply_template(slide, slide_number)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(10.33), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.name = "Times New Roman"
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Content
    if bullet_points:
        content_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.0), Inches(9.0), Inches(4.5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf_content.paragraphs[0]
            else:
                p = tf_content.add_paragraph()
            p.text = point
            p.font.size = Pt(24)
            p.font.name = "Times New Roman"
            p.font.color.rgb = RGBColor(0, 0, 0)
            if not point.startswith(" -") and not point.startswith("("):
                p.level = 0
            else:
                p.level = 1
                p.font.size = Pt(20)
                
    return slide

# Slide 1: Project Title
slide1_layout = prs.slide_layouts[6] # Blank
slide1 = prs.slides.add_slide(slide1_layout)
apply_template(slide1, 1)

# Title
title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(10.33), Inches(1.0))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "ShieldX"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(54)
p.font.name = "Times New Roman"

# Team Members Table
table_shape = slide1.shapes.add_table(1, 3, Inches(2.5), Inches(2.5), Inches(8.33), Inches(0.6))
table = table_shape.table
table.cell(0,0).text = "Jogonnath Das Talukder"
table.cell(0,1).text = "Shrestta Das"
table.cell(0,2).text = "Khandoker Syed Shovon"
for cell in table.iter_cells():
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(220, 220, 220)
    # Add thin border
    # cell borders are a bit complex in python-pptx, so we'll just rely on background color
    for par in cell.text_frame.paragraphs:
        par.alignment = PP_ALIGN.CENTER
        par.font.name = "Times New Roman"
        par.font.size = Pt(18)
        par.font.color.rgb = RGBColor(0,0,0)

sup_box = slide1.shapes.add_textbox(Inches(2.0), Inches(3.5), Inches(9.33), Inches(1.0))
tf3 = sup_box.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Supervisor: Shahriar Arefin Zummon\nLecturer, Dept. of CSE, Leading University"
p3.alignment = PP_ALIGN.CENTER
p3.font.size = Pt(24)
p3.font.name = "Times New Roman"

desc_box = slide1.shapes.add_textbox(Inches(2.5), Inches(5.5), Inches(8.33), Inches(1.5))
tf4 = desc_box.text_frame
tf4.word_wrap = True
p4 = tf4.paragraphs[0]
p4.text = "This slide contains the official project title, the names of all team members, and the name of the project supervisor. The title is clear, concise, and professional."
p4.alignment = PP_ALIGN.JUSTIFY
p4.font.size = Pt(20)
p4.font.name = "Times New Roman"
p4.font.italic = True

# Slide 2: Presentation Outline
add_styled_slide("Presentation Outline", [
    "• Team Overview",
    "• Project Overview & Goal",
    "• Functional Requirements",
    "• Non-Functional Requirements",
    "• Technical Diagrams",
    "• Project Demonstration",
    "• Discussion / Conclusion",
    "• Demo"
], 2)

# Slide 3: Team Overview
add_styled_slide("Team Overview*", [
    "• Jogonnath Das Talukder",
    " - Developed the frontend API and Firebase integration",
    "• Shrestta Das",
    " - Implemented SOS systems, mappings, and Supabase integration",
    "• Khandoker Syed Shovon",
    " - Designed the database schema and implemented admin dashboard"
], 3)

# Slide 4: Project Overview & Goal
add_styled_slide("Project Overview & Goal", [
    "• Description of the project:",
    " - A comprehensive mobile application bridging the gap between citizens and law enforcement.",
    "• Main objectives and goals:",
    " - Empower citizens to report crimes, track complaints, and trigger SOS alerts.",
    " - Provide admins with a command dashboard to manage reports efficiently.",
    "• Problem statement & motivation:",
    " - The lack of a centralized, real-time emergency reporting and tracking system within the Sylhet Metropolitan Police jurisdiction."
], 4)

# Slide 5: Functional Requirements
add_styled_slide("Functional Requirements", [
    "• List of functional requirements:",
    " - Secure Onboarding (OTP & NID Verification)",
    " - Smart Complaint Submission",
    " - Real-time Tracking & Interactive Police Map",
    " - SOS Emergency Broadcasting",
    " - Admin Command Dashboard",
    "• Use Case Diagram:",
    " - (Insert Use Case Diagram Here)"
], 5)

# Slide 6: Non-Functional Requirements
add_styled_slide("Non-Functional Requirements", [
    "• Operating system environment:",
    " - Android & iOS (Cross-platform)",
    "• Frameworks and tools:",
    " - Flutter (Dart) for UI",
    " - Supabase (PostgreSQL) for Backend",
    " - Riverpod for State Management",
    " - GoRouter for Routing",
    " - flutter_map for Geolocation Services"
], 6)

# Slide 7: Technical diagrams
add_styled_slide("Technical diagrams", [
    "• Technical diagrams (ER Diagram / DFD / Class Diagram)",
    " - (Please insert the diagrams on this and the following slides)"
], 7)

# Slide 8: Project Demonstration
add_styled_slide("Project Demonstration", [
    "Insert video or slideshow of project demo",
    "• 1-minute overview covering key features:"
], 8)

# Slide 9: Discussion / Conclusion
add_styled_slide("Discussion / Conclusion", [
    "• Limitations:",
    " - Reliance on active internet and GPS.",
    "• Restrictions:",
    " - Currently scoped to SMP jurisdiction.",
    "• Future scope and improvements:",
    " - Scale up to a national platform.",
    " - Incorporate AI for crime pattern prediction.",
    " - Implement multilingual interfaces."
], 9)

# Slide 10: Demo
add_styled_slide("Demo", [
    "• Present the Functional Demo"
], 10)

prs.save("ShieldX_Template_Matched_Presentation.pptx")
print("Presentation saved to ShieldX_Template_Matched_Presentation.pptx")
