import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

def add_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
    return slide

# Slide 1: Project Title
title_slide_layout = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(title_slide_layout)
slide1.shapes.title.text = "ShieldX\nCrime Reporting & Emergency Management System"
slide1.placeholders[1].text = (
    "Team Members:\n"
    "Jogonnath Das Talukder (0182320012101060)\n"
    "Shrestta Das (0182320012101273)\n"
    "Khandoker Syed Shovon (0182320012101400)\n\n"
    "Supervisor:\n"
    "Shahriar Arefin Zummon\n"
    "Lecturer, Dept. of CSE, Leading University"
)

# Slide 2: Presentation Outline
add_slide("Presentation Outline", [
    "Team Overview",
    "Project Overview & Goal",
    "Functional Requirements",
    "Non-Functional Requirements",
    "Technical Diagrams",
    "Project Demonstration",
    "Discussion / Conclusion",
    "Demo"
])

# Slide 3: Team Overview
add_slide("Team Overview", [
    "Jogonnath Das Talukder (0182320012101060): Contribution details",
    "Shrestta Das (0182320012101273): Contribution details",
    "Khandoker Syed Shovon (0182320012101400): Contribution details"
])

# Slide 4: Project Overview & Goal
add_slide("Project Overview & Goal", [
    "Description: ShieldX is a mobile app to bridge the gap between citizens and law enforcement in Sylhet Metropolitan Police (SMP) jurisdiction.",
    "Main Objectives: Empower citizens to report crimes, track complaints in real-time, trigger SOS alerts, and provide admins with a command dashboard.",
    "Problem Statement: Lack of efficient crime reporting, real-time tracking, and immediate emergency response systems."
])

# Slide 5: Functional Requirements
add_slide("Functional Requirements", [
    "Secure Onboarding (OTP & NID Verification)",
    "Smart Complaint Submission (Form, Media Upload, Geolocation)",
    "Real-time Tracking of Complaint Journey",
    "SOS Emergency Panic Button with Live Geolocation",
    "Admin Command Dashboard & Complaint Management",
    "(Insert Use Case Diagram Here)"
])

# Slide 6: Non-Functional Requirements
add_slide("Non-Functional Requirements", [
    "Operating system environment:",
    " - Android & iOS (Cross-platform)",
    "Frameworks and tools:",
    " - Flutter (Dart) for UI Framework",
    " - Supabase (PostgreSQL) for Backend & Database",
    " - Riverpod for State Management",
    " - GoRouter for Routing",
    " - flutter_map + OpenStreetMap for Mapping Services"
])

# Slide 7: Technical diagrams
add_slide("Technical diagrams", [
    "Entity-Relationship (ER) Diagram",
    "Data Flow Diagram (DFD Level 0 or 1)",
    "System Architecture Diagram",
    "(Insert Technical Diagrams Here)"
])

# Slide 8: Project Demonstration
add_slide("Project Demonstration", [
    "Insert video or slideshow of project demo",
    "1-minute overview covering:",
    " - User Registration and Authentication",
    " - Creating a Complaint and Live Tracking",
    " - Admin Dashboard Insights",
    " - Sending an SOS Alert"
])

# Slide 9: Discussion / Conclusion
add_slide("Discussion / Conclusion", [
    "Limitations: Relies on active internet and GPS availability.",
    "Restrictions: Initially tailored for SMP jurisdiction.",
    "Future scope and improvements:",
    " - Scale up to a national platform.",
    " - Incorporate AI for crime pattern prediction.",
    " - Implement multilingual interfaces."
])

# Slide 10: Demo
add_slide("Demo", [
    "Present the Functional Demo"
])

prs.save("ShieldX_Defense_Presentation.pptx")
print("Presentation generated successfully at ShieldX_Defense_Presentation.pptx")
